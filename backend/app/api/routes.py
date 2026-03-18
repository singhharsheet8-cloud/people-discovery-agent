import asyncio
import csv
import io
import json
import time
import uuid
import logging
from datetime import datetime, timezone
from fastapi import APIRouter, Query, HTTPException, Depends, Body
from fastapi.responses import JSONResponse, StreamingResponse
from pydantic import BaseModel, Field, field_validator
from sqlalchemy import select, func, desc, asc
from sqlalchemy.ext.asyncio import AsyncSession
from app.db import get_session_factory
from app.models.db_models import Person, PersonSource, DiscoveryJob, PersonVersion, AdminUser, compute_name_key
from app.cache import cleanup_expired_cache
from app.auth import require_admin, require_api, require_viewer
from app.embeddings import update_person_embedding
from passlib.hash import bcrypt
from app.rate_limiter import rate_limiter
from app.config import get_settings
from app.intelligence import (
    analyze_sentiment,
    map_relationships,
    calculate_influence_score,
    generate_meeting_prep,
    verify_facts,
)

logger = logging.getLogger(__name__)
router = APIRouter(prefix="/api")


async def _set_job_step(job_id: str, step: str) -> None:
    """Write the current pipeline step to the DB so the frontend can poll it."""
    try:
        factory = get_session_factory()
        async with factory() as session:
            job = (await session.execute(
                select(DiscoveryJob).where(DiscoveryJob.id == job_id)
            )).scalar_one_or_none()
            if job:
                job.current_step = step
                await session.commit()
    except Exception:
        pass  # Never crash the pipeline over a step-tracking failure


# --- Request/Response models ---

class DiscoverRequest(BaseModel):
    name: str = Field(..., min_length=1, max_length=255)
    company: str = Field("", max_length=255)
    role: str = Field("", max_length=255)
    location: str = Field("", max_length=255)
    linkedin_url: str = Field("", max_length=500)
    twitter_handle: str = Field("", max_length=100)
    github_username: str = Field("", max_length=100)
    instagram_handle: str = Field("", max_length=100)
    context: str = Field("", max_length=2000)

    @field_validator("linkedin_url")
    @classmethod
    def validate_linkedin_url(cls, v: str) -> str:
        if v and not v.startswith(("https://linkedin.com/", "https://www.linkedin.com/", "http://linkedin.com/", "http://www.linkedin.com/")):
            raise ValueError("LinkedIn URL must start with https://linkedin.com/ or https://www.linkedin.com/")
        return v

    @field_validator("twitter_handle")
    @classmethod
    def validate_twitter_handle(cls, v: str) -> str:
        if v:
            v = v.lstrip("@")
            if not v.replace("_", "").isalnum():
                raise ValueError("Twitter handle must contain only alphanumeric characters and underscores")
        return v

    @field_validator("instagram_handle")
    @classmethod
    def validate_instagram_handle(cls, v: str) -> str:
        if v:
            v = v.lstrip("@")
            if not v.replace("_", "").replace(".", "").isalnum():
                raise ValueError("Instagram handle must contain only alphanumeric characters, underscores, and periods")
        return v

    @field_validator("github_username")
    @classmethod
    def validate_github_username(cls, v: str) -> str:
        if v and not v.replace("-", "").replace("_", "").isalnum():
            raise ValueError("GitHub username must contain only alphanumeric characters, hyphens, and underscores")
        return v


class DiscoverResponse(BaseModel):
    job_id: str
    status: str
    message: str


class PersonSummary(BaseModel):
    id: str
    name: str
    company: str | None = None
    current_role: str | None = None
    image_url: str | None = None
    confidence_score: float = 0.0
    status: str = "discovered"
    sources_count: int = 0
    created_at: str = ""
    updated_at: str = ""


class PersonDetail(BaseModel):
    id: str
    name: str
    current_role: str | None = None
    company: str | None = None
    location: str | None = None
    bio: str | None = None
    image_url: str | None = None
    education: list | None = None
    key_facts: list | None = None
    social_links: dict | None = None
    expertise: list | None = None
    notable_work: list | None = None
    career_timeline: list | None = None
    skills: list | None = None
    projects: list | None = None
    recommendations: list | None = None
    languages: list | None = None
    followers_count: int | None = None
    blog_url: str | None = None
    confidence_score: float = 0.0
    reputation_score: float | None = None
    influence_score: float | None = None
    status: str = "discovered"
    version: int = 1
    sources: list[dict] = []
    jobs: list[dict] = []
    created_at: str = ""
    updated_at: str = ""


class PersonUpdate(BaseModel):
    name: str | None = Field(None, max_length=255)
    current_role: str | None = Field(None, max_length=500)
    company: str | None = Field(None, max_length=500)
    location: str | None = Field(None, max_length=255)
    bio: str | None = Field(None, max_length=20000)
    image_url: str | None = None
    # Extended fields for manual data enrichment
    confidence_score: float | None = None
    reputation_score: float | None = None
    status: str | None = None
    key_facts: list | None = None
    education: list | None = None
    career_timeline: list | None = None
    expertise: list | None = None
    notable_work: list | None = None
    social_links: dict | None = None
    skills: list | None = None
    projects: list | None = None
    recommendations: list | None = None
    followers_count: int | None = None
    blog_url: str | None = Field(None, max_length=500)


class JobDetail(BaseModel):
    id: str
    person_id: str | None = None
    status: str
    input_params: dict = {}
    cost_breakdown: dict | None = None
    total_cost: float = 0.0
    latency_ms: float | None = None
    sources_hit: int = 0
    cache_hits: int = 0
    error_message: str | None = None
    created_at: str = ""
    completed_at: str | None = None


# --- Discovery endpoint ---

@router.post("/discover", response_model=DiscoverResponse)
async def discover_person(request: DiscoverRequest, _auth=Depends(require_api)):
    """Single-shot person discovery. Returns a job ID for polling."""
    settings = get_settings()
    factory = get_session_factory()

    async with factory() as session:
        running_count = (await session.execute(
            select(func.count()).select_from(DiscoveryJob).where(DiscoveryJob.status == "running")
        )).scalar() or 0
        if running_count >= settings.max_concurrent_jobs:
            raise HTTPException(
                status_code=429,
                detail=f"Too many concurrent discovery jobs ({running_count}/{settings.max_concurrent_jobs}). Try again shortly.",
            )

        today_start = datetime.now(timezone.utc).replace(hour=0, minute=0, second=0, microsecond=0)
        daily_count = (await session.execute(
            select(func.count()).select_from(DiscoveryJob).where(DiscoveryJob.created_at >= today_start)
        )).scalar() or 0
        if daily_count >= settings.max_daily_discoveries:
            raise HTTPException(
                status_code=429,
                detail=f"Daily discovery limit reached ({settings.max_daily_discoveries}). Try again tomorrow.",
            )

    job_id = str(uuid.uuid4())
    input_data = request.model_dump(exclude_defaults=False)

    async with factory() as session:
        job = DiscoveryJob(
            id=job_id,
            input_params=json.dumps(input_data),
            status="running",
        )
        session.add(job)
        await session.commit()

    asyncio.create_task(_run_discovery(job_id, input_data))

    return DiscoverResponse(
        job_id=job_id,
        status="running",
        message="Discovery started. Poll GET /api/jobs/{job_id} for status.",
    )


class BatchDiscoverRequest(BaseModel):
    persons: list[DiscoverRequest] = Field(..., min_length=1, max_length=50)


@router.post("/discover/batch")
async def batch_discover(request: BatchDiscoverRequest, _auth=Depends(require_api)):
    """Start discovery for multiple people. Returns list of job IDs."""
    settings = get_settings()
    factory = get_session_factory()

    async with factory() as session:
        running_count = (await session.execute(
            select(func.count()).select_from(DiscoveryJob).where(DiscoveryJob.status == "running")
        )).scalar() or 0
        if running_count + len(request.persons) > settings.max_concurrent_jobs:
            raise HTTPException(
                status_code=429,
                detail=f"Too many concurrent jobs. {running_count} running, {len(request.persons)} requested, max {settings.max_concurrent_jobs}",
            )

    jobs = []
    for person in request.persons:
        job_id = str(uuid.uuid4())
        input_data = person.model_dump(exclude_defaults=False)

        async with factory() as session:
            job = DiscoveryJob(id=job_id, input_params=json.dumps(input_data), status="running")
            session.add(job)
            await session.commit()

        asyncio.create_task(_run_discovery(job_id, input_data))
        jobs.append({"job_id": job_id, "name": person.name, "status": "running"})

    return {"jobs": jobs, "total": len(jobs)}


def _normalize_name(name: str) -> str:
    """Lowercase and strip punctuation from a name string."""
    import re as _re
    return _re.sub(r"[^\w\s]", "", name.strip().lower())


def _names_match(a: str, b: str) -> bool:
    """Return True if names refer to the same person (case-insensitive, word-order invariant)."""
    if not a or not b:
        return False
    norm_a = set(_normalize_name(a).split())
    norm_b = set(_normalize_name(b).split())
    if not norm_a or not norm_b:
        return False
    return norm_a == norm_b


def _companies_match(a: str | None, b: str | None) -> bool:
    """Return True if two company strings refer to the same company."""
    _SUFFIXES = {"inc", "corp", "ltd", "llc", "co", "corporation", "limited", "incorporated"}
    if a is None and b is None:
        return True
    if a is None or b is None:
        return False
    def _clean(s: str) -> str:
        import re as _re
        s = s.strip().lower()
        parts = s.split()
        parts = [p.rstrip(".,") for p in parts if p.rstrip(".,") not in _SUFFIXES]
        return " ".join(parts)
    return _clean(a) == _clean(b)


async def _find_existing_person(session: AsyncSession, name: str, company: str | None) -> Person | None:
    """
    Return an existing Person whose canonical name_key matches the incoming name.

    This is the single source of truth for duplicate detection. The name_key is
    computed by compute_name_key() (aliases collapsed, middle names dropped, words
    sorted) and enforced at the DB level by a UNIQUE constraint + trigger, so even
    direct SQL writes cannot create duplicates.

    We do a single indexed lookup on name_key — O(1), no full-table scans.
    """
    key = compute_name_key(name)
    if not key:
        return None
    row = (await session.execute(
        select(Person).where(Person.name_key == key)
    )).scalar_one_or_none()
    return row


def _merge_scalar(old_val: str | None, new_val: str | None) -> str | None:
    """Pick the richer scalar value (longer non-empty string wins)."""
    if not new_val:
        return old_val
    if not old_val:
        return new_val
    return new_val if len(new_val) >= len(old_val) else old_val


def _merge_list_field(old_list: list | None, new_list: list | None) -> list:
    """Union two lists, deduplicating by normalized content."""
    old_list = old_list or []
    new_list = new_list or []
    seen: set[str] = set()
    merged: list = []
    for item in old_list + new_list:
        if isinstance(item, str):
            key = item.strip().lower()
        elif isinstance(item, dict):
            # For timeline/career entries, dedup by type+title+company
            if "title" in item or "company" in item:
                key = "|".join([
                    (item.get("type") or "").lower().strip(),
                    (item.get("title") or item.get("description", "")).lower().strip(),
                    (item.get("company") or "").lower().strip(),
                ])
            else:
                key = json.dumps(item, sort_keys=True)
        else:
            key = str(item)
        if key not in seen:
            seen.add(key)
            merged.append(item)
    return merged


def _merge_dict_field(old_dict: dict | None, new_dict: dict | None) -> dict:
    """Merge two dicts, preferring non-null values from the newer one."""
    old_dict = old_dict or {}
    new_dict = new_dict or {}
    merged = {**old_dict}
    for k, v in new_dict.items():
        if v is not None and (v != "" or k not in merged):
            merged[k] = v
    return merged


async def _run_discovery(job_id: str, input_data: dict):
    """Run the LangGraph agent, then merge results into existing or new person."""
    from app.agent.graph import graph

    start_time = time.time()
    factory = get_session_factory()

    try:
        initial_state = {
            "messages": [],
            "input": input_data,
            "search_queries": [],
            "search_results": [],
            "analyzed_results": {},
            "enrichment": {},
            "sentiment": {},
            "confidence_score": 0.0,
            "person_profile": None,
            "cost_tracker": {},
            "status": "started",
            # Iterative enrichment / disambiguation fields
            "iteration": 0,
            "identity_anchors": [],
            "filtered_results": [],
            "refinement_queries": [],
            "refinement_signals": [],
            "executed_query_hashes": [],
            "abort_reason": None,
        }

        # Stream node completions so we can write current_step to the DB in real time.
        # LangGraph's astream_events yields events per node; we update the job row
        # after each node completes so the frontend polling sees live progress.
        _STEP_LABELS: dict[str, str] = {
            "plan_searches":            "planning",
            "execute_searches":         "searching",
            "disambiguate":             "disambiguating",
            "filter_results":           "filtering",
            "analyze_results":          "analyzing",
            "enrich_data":              "enriching",
            "iterative_enrich":         "iterating",
            "generate_targeted_queries":"refining",
            "synthesize_profile":       "synthesizing",   # sentiment runs concurrently inside this node
            "verify_profile":           "verifying",
        }
        result = None
        async for event in graph.astream_events(initial_state, version="v2"):
            kind = event.get("event", "")
            node = event.get("name", "")
            if kind == "on_chain_start" and node in _STEP_LABELS:
                await _set_job_step(job_id, _STEP_LABELS[node])
            elif kind == "on_chain_end" and node == "LangGraph":
                result = event.get("data", {}).get("output", {})

        if result is None:
            result = {}
        profile = result.get("person_profile", {})

        _MIN_SAVE_CONFIDENCE = 0.35
        profile_confidence = float(profile.get("confidence_score", 0))
        if profile.get("abort_reason") or profile_confidence < _MIN_SAVE_CONFIDENCE:
            logger.warning(
                "Profile for %s rejected (confidence=%.3f, abort=%s) — not saving to DB",
                input_data.get("name"), profile_confidence, profile.get("abort_reason"),
            )
            async with factory() as session:
                job_row = (await session.execute(
                    select(DiscoveryJob).where(DiscoveryJob.id == job_id)
                )).scalar_one_or_none()
                if job_row:
                    job_row.status = "failed"
                    job_row.completed_at = datetime.now(timezone.utc)
                    job_row.latency_ms = (time.time() - start_time) * 1000
                await session.commit()
            return

        sentiment = result.get("sentiment", {})
        if sentiment:
            profile["sentiment"] = sentiment
        enrichment = result.get("enrichment", {})
        if enrichment.get("image_url") and not profile.get("image_url"):
            profile["image_url"] = enrichment["image_url"]
        elapsed_ms = (time.time() - start_time) * 1000

        async with factory() as session:
            job_row = (await session.execute(
                select(DiscoveryJob).where(DiscoveryJob.id == job_id)
            )).scalar_one_or_none()

            existing_person: Person | None = None

            if job_row and job_row.person_id:
                existing_person = (await session.execute(
                    select(Person).where(Person.id == job_row.person_id)
                )).scalar_one_or_none()

            if not existing_person:
                new_name = profile.get("name", input_data.get("name", ""))
                new_company = profile.get("company", input_data.get("company"))
                existing_person = await _find_existing_person(session, new_name, new_company)

            is_merge = existing_person is not None

            if is_merge:
                person = existing_person
                logger.info(f"Merging discovery into existing person {person.id} ({person.name})")

                person.name = profile.get("name") or person.name
                person.current_role = _merge_scalar(person.current_role, profile.get("current_role"))
                person.company = _merge_scalar(person.company, profile.get("company"))
                person.location = _merge_scalar(person.location, profile.get("location"))
                person.bio = _merge_scalar(person.bio, profile.get("bio"))
                new_conf = profile.get("confidence_score", result.get("confidence_score", 0))
                person.confidence_score = round(
                    (person.confidence_score + new_conf) / 2, 3
                )
                if profile.get("image_url") and not person.image_url:
                    person.image_url = profile["image_url"]
                person.status = "discovered"

                list_fields = ("education", "key_facts", "expertise", "notable_work", "career_timeline",
                               "skills", "projects", "recommendations")
                for field in list_fields:
                    old_val = person.get_json(field)
                    new_val = profile.get(field)
                    merged = _merge_list_field(old_val, new_val)
                    if merged:
                        person.set_json(field, merged)

                dict_fields = ("social_links",)
                for field in dict_fields:
                    old_val = person.get_json(field)
                    new_val = profile.get(field)
                    merged = _merge_dict_field(old_val, new_val)
                    if merged:
                        person.set_json(field, merged)

                # Scalar enrichment fields — only update if new value is provided and current is empty
                if profile.get("followers_count") and not person.followers_count:
                    person.followers_count = profile["followers_count"]
                if profile.get("blog_url") and not person.blog_url:
                    person.blog_url = profile["blog_url"]

                # Auto-populate social_links.website from personal_website sources (merge path)
                for sr in result.get("search_results", []):
                    if isinstance(sr, dict) and sr.get("source_type") == "personal_website":
                        u = sr.get("url", "").split("?")[0].rstrip("/")
                        if u and "linkedin.com" not in u:
                            sl = person.get_json("social_links") or {}
                            if not sl.get("website"):
                                sl["website"] = u
                                person.set_json("social_links", sl)
                            if not person.blog_url:
                                person.blog_url = u
                            break

                person.version += 1
                person.updated_at = datetime.now(timezone.utc)

            else:
                person = Person(
                    name=profile.get("name", input_data.get("name", "Unknown")),
                    current_role=profile.get("current_role"),
                    company=profile.get("company"),
                    location=profile.get("location"),
                    bio=profile.get("bio"),
                    image_url=profile.get("image_url"),
                    confidence_score=profile.get("confidence_score", result.get("confidence_score", 0)),
                    status="discovered",
                )
                for field in ("education", "key_facts", "social_links", "expertise", "notable_work",
                              "career_timeline", "skills", "projects", "recommendations"):
                    val = profile.get(field)
                    if val:
                        person.set_json(field, val)
                if profile.get("followers_count"):
                    person.followers_count = profile["followers_count"]
                if profile.get("blog_url"):
                    person.blog_url = profile["blog_url"]
                session.add(person)

            # ── Auto-populate social_links.website from personal_website sources ──
            # Mirror pattern: LinkedIn URL goes into social_links.linkedin, so
            # personal_website URL should go into social_links.website and blog_url.
            personal_site_url = None
            for sr in result.get("search_results", []):
                if isinstance(sr, dict) and sr.get("source_type") == "personal_website":
                    u = sr.get("url", "").split("?")[0].rstrip("/")
                    if u and "linkedin.com" not in u:
                        personal_site_url = u
                        break
            if personal_site_url:
                # Merge into social_links.website
                sl = person.get_json("social_links") or {}
                if not sl.get("website"):
                    sl["website"] = personal_site_url
                    person.set_json("social_links", sl)
                # Also populate blog_url if empty
                if not person.blog_url:
                    person.blog_url = personal_site_url

            await session.flush()

            existing_keys: set[tuple[str, str]] = set()
            if is_merge:
                existing_sources = (await session.execute(
                    select(PersonSource.url, PersonSource.source_type).where(PersonSource.person_id == person.id)
                )).all()
                existing_keys = {(row[0] or "", row[1] or "") for row in existing_sources}

            _SOURCE_RELEVANCE_FLOOR = 0.45
            new_keys: set[tuple[str, str]] = set()
            for source in profile.get("sources", []):
                url = source.get("url", "")
                stype = source.get("platform", "web")
                rel = source.get("relevance_score", source.get("confidence", 0.5))
                if float(rel) < _SOURCE_RELEVANCE_FLOOR:
                    continue
                key = (url, stype)
                if key in existing_keys or key in new_keys:
                    continue
                new_keys.add(key)
                ps = PersonSource(
                    person_id=person.id,
                    source_type=stype,
                    platform=stype,
                    url=url,
                    title=source.get("title", ""),
                    raw_content=source.get("snippet", ""),
                    relevance_score=float(rel),
                    source_reliability=source.get("confidence", _get_source_reliability(stype)),
                )
                session.add(ps)

            for sr in result.get("search_results", []):
                if isinstance(sr, dict):
                    url = sr.get("url", "")
                    stype = sr.get("source_type", "web")
                    rel_score = float(sr.get("relevance_score", sr.get("score", 0.5)))
                    if rel_score < _SOURCE_RELEVANCE_FLOOR:
                        continue
                    key = (url, stype)
                    if key in existing_keys or key in new_keys:
                        continue
                    new_keys.add(key)
                    src_rel = sr.get("source_reliability", _get_source_reliability(stype))
                    ps = PersonSource(
                        person_id=person.id,
                        source_type=stype,
                        platform=stype,
                        url=url,
                        title=sr.get("title", ""),
                        raw_content=sr.get("content", "")[:5000],
                        structured_data=json.dumps(sr.get("structured")) if sr.get("structured") else None,
                        relevance_score=rel_score,
                        source_reliability=float(src_rel),
                        scorer_reason=(sr.get("scorer_reason") or "")[:200] or None,
                    )
                    session.add(ps)

            trigger = "re-search" if is_merge else "initial"
            version = PersonVersion(
                person_id=person.id,
                version_number=person.version,
                profile_snapshot=json.dumps(profile),
                trigger=trigger,
            )
            session.add(version)

            if job_row:
                job_row.person_id = person.id
                job_row.status = "completed"
                job_row.total_cost = result.get("cost_tracker", {}).get("total", 0.0)
                job_row.cost_breakdown = json.dumps(result.get("cost_tracker", {}))
                job_row.latency_ms = elapsed_ms
                all_sr = result.get("search_results", [])
                job_row.sources_hit = sum(
                    1 for s in all_sr
                    if isinstance(s, dict) and float(s.get("relevance_score", s.get("confidence", 0))) >= _SOURCE_RELEVANCE_FLOOR
                )
                job_row.completed_at = datetime.now(timezone.utc)

            await session.commit()
            action = "merged into" if is_merge else "created"
            logger.info(
                f"Discovery complete for job {job_id}: {action} {person.name} "
                f"(v{person.version}, {len(new_keys)} new sources, {elapsed_ms:.0f}ms)"
            )

            # Generate and persist semantic embedding (non-blocking; errors are logged)
            await update_person_embedding(session, person)
            await session.commit()

        # Post-save image fill — if the agent didn't find an image, run the
        # resolver now with the fully known LinkedIn handle so Tier 2/3
        # (SerpAPI Google Images) can fill it in.
        if not person.image_url:
            person_id_for_img = person.id
            person_name_for_img = person.name
            person_company_for_img = person.company
            search_results_for_img = result.get("search_results", [])
            asyncio.create_task(
                _fill_image_post_save(
                    person_id_for_img,
                    person_name_for_img,
                    person_company_for_img,
                    search_results_for_img,
                )
            )

        from app.api.webhooks import fire_webhooks
        event = "person.updated" if is_merge else "job.completed"
        await fire_webhooks(event, {
            "job_id": job_id,
            "person_id": person.id,
            "person_name": person.name,
            "status": "completed",
            "merged": is_merge,
            "total_cost": result.get("cost_tracker", {}).get("total", 0.0),
            "latency_ms": round(elapsed_ms),
            "sources_hit": len(result.get("search_results", [])),
            "new_sources_added": len(new_keys),
        })

    except Exception as e:
        logger.error(f"Discovery failed for job {job_id}: {e}")
        async with factory() as session:
            job_stmt = select(DiscoveryJob).where(DiscoveryJob.id == job_id)
            job_result = await session.execute(job_stmt)
            job = job_result.scalar_one_or_none()
            if job:
                job.status = "failed"
                job.error_message = str(e)[:2000]
                job.latency_ms = (time.time() - start_time) * 1000
                job.completed_at = datetime.now(timezone.utc)
            await session.commit()


async def _fill_image_post_save(
    person_id: str,
    name: str,
    company: str | None,
    search_results: list[dict],
) -> None:
    """
    Background task: resolve the profile image after the person is already
    saved, then write it back to the DB.  Called only when the main agent
    run produced no image_url.
    """
    try:
        from app.tools.image_resolver import resolve_profile_image
        image_url = await resolve_profile_image(
            name=name, company=company, search_results=search_results
        )
        if not image_url:
            return
        factory = get_session_factory()
        async with factory() as session:
            person = (await session.execute(
                select(Person).where(Person.id == person_id)
            )).scalar_one_or_none()
            if person and not person.image_url:
                person.image_url = image_url
                await session.commit()
                logger.info(f"[image] post-save fill: {name!r} → {image_url[:80]}")
    except Exception as e:
        logger.warning(f"[image] post-save fill failed for {name!r}: {e}")


def _get_source_reliability(platform: str) -> float:
    return {
        "linkedin_profile": 0.95, "linkedin_posts": 0.85,
        "github": 0.9, "twitter": 0.7, "youtube_transcript": 0.85,
        "news": 0.8, "academic": 0.95, "scholar": 0.95,
        "web": 0.6, "reddit": 0.4, "medium": 0.7,
        "crunchbase": 0.9, "instagram": 0.5, "firecrawl": 0.7,
    }.get(platform, 0.5)


# --- Jobs endpoint ---

def _validate_uuid(value: str, label: str = "ID") -> str:
    try:
        uuid.UUID(value)
    except ValueError:
        raise HTTPException(status_code=400, detail=f"Invalid {label} format")
    return value


@router.get("/jobs/{job_id}")
async def get_job(job_id: str, _auth=Depends(require_api)):
    _validate_uuid(job_id, "job_id")
    factory = get_session_factory()
    async with factory() as session:
        stmt = select(DiscoveryJob).where(DiscoveryJob.id == job_id)
        result = await session.execute(stmt)
        job = result.scalar_one_or_none()
        if not job:
            raise HTTPException(status_code=404, detail="Job not found")

        resp = {
            "id": job.id,
            "person_id": job.person_id,
            "status": job.status,
            "current_step": job.current_step,
            "input_params": json.loads(job.input_params) if job.input_params else {},
            "cost_breakdown": json.loads(job.cost_breakdown) if job.cost_breakdown else None,
            "total_cost": job.total_cost,
            "latency_ms": job.latency_ms,
            "sources_hit": job.sources_hit,
            "cache_hits": job.cache_hits,
            "error_message": job.error_message,
            "created_at": job.created_at.isoformat() if job.created_at else "",
            "completed_at": job.completed_at.isoformat() if job.completed_at else None,
        }

        # If completed, include the person profile
        if job.status == "completed" and job.person_id:
            person_stmt = select(Person).where(Person.id == job.person_id)
            person_result = await session.execute(person_stmt)
            person = person_result.scalar_one_or_none()
            if person:
                resp["profile"] = _person_to_dict(person)

        return resp


# --- Persons CRUD ---


class SemanticSearchRequest(BaseModel):
    query: str = Field(..., min_length=2, max_length=500, description="Natural language query")
    limit: int = Field(10, ge=1, le=50)
    min_similarity: float = Field(0.0, ge=0.0, le=1.0, description="Minimum cosine similarity (0–1)")


async def _run_semantic_search(q: str, limit: int, min_similarity: float):
    from app.embeddings import semantic_search

    factory = get_session_factory()
    async with factory() as session:
        results = await semantic_search(
            session, query=q, limit=limit, min_similarity=min_similarity
        )
    return {"results": results, "count": len(results), "query": q}


@router.post("/persons/semantic-search")
async def semantic_search_persons_post(
    request: SemanticSearchRequest,
    _auth=Depends(require_api),
):
    """Semantic search via POST body — type your query in JSON."""
    return await _run_semantic_search(request.query, request.limit, request.min_similarity)


@router.get("/persons/semantic-search")
async def semantic_search_persons(
    q: str = Query(..., min_length=2, max_length=500, description="Natural language query"),
    limit: int = Query(10, ge=1, le=50),
    min_similarity: float = Query(0.0, ge=0.0, le=1.0, description="Minimum cosine similarity (0–1)"),
    _auth=Depends(require_api),
):
    """Semantic search via query params (GET)."""
    return await _run_semantic_search(q, limit, min_similarity)


@router.get("/persons")
async def list_persons(
    page: int = Query(1, ge=1),
    per_page: int = Query(20, ge=1, le=100),
    search: str = Query("", max_length=255),
    company: str = Query("", max_length=255),
    location: str = Query("", max_length=255),
    min_confidence: float = Query(0.0, ge=0.0, le=1.0),
    sort_by: str = Query("updated_at", pattern="^(updated_at|name|confidence_score|created_at)$"),
    sort_order: str = Query("desc", pattern="^(asc|desc)$"),
    _auth=Depends(require_api),
):
    factory = get_session_factory()
    async with factory() as session:
        sort_col = getattr(Person, sort_by, Person.updated_at)
        order_fn = desc if sort_order == "desc" else asc
        query = select(Person).order_by(order_fn(sort_col))

        if search:
            safe = search.replace("%", r"\%").replace("_", r"\_")
            query = query.where(
                (Person.name.ilike(f"%{safe}%")) | (Person.company.ilike(f"%{safe}%"))
                | (Person.current_role.ilike(f"%{safe}%"))
            )

        if company:
            safe_c = company.replace("%", r"\%").replace("_", r"\_")
            query = query.where(Person.company.ilike(f"%{safe_c}%"))

        if location:
            safe_l = location.replace("%", r"\%").replace("_", r"\_")
            query = query.where(Person.location.ilike(f"%{safe_l}%"))

        if min_confidence > 0:
            query = query.where(Person.confidence_score >= min_confidence)

        # Count
        count_q = select(func.count()).select_from(query.subquery())
        total = (await session.execute(count_q)).scalar() or 0

        # Paginate
        query = query.offset((page - 1) * per_page).limit(per_page)
        result = await session.execute(query)
        persons = result.scalars().all()

        # Get source counts
        items = []
        for p in persons:
            src_count_stmt = select(func.count()).select_from(PersonSource).where(PersonSource.person_id == p.id)
            src_count = (await session.execute(src_count_stmt)).scalar() or 0
            items.append({
                "id": p.id,
                "name": p.name,
                "company": p.company,
                "current_role": p.current_role,
                "image_url": p.image_url,
                "confidence_score": p.confidence_score,
                "status": p.status,
                "sources_count": src_count,
                "created_at": p.created_at.isoformat() if p.created_at else "",
                "updated_at": p.updated_at.isoformat() if p.updated_at else "",
            })

        return {"items": items, "total": total, "page": page, "per_page": per_page}


@router.get("/persons/{person_id}")
async def get_person(person_id: str, _auth=Depends(require_api)):
    _validate_uuid(person_id, "person_id")
    factory = get_session_factory()
    async with factory() as session:
        person = (await session.execute(
            select(Person).where(Person.id == person_id)
        )).scalar_one_or_none()
        if not person:
            raise HTTPException(status_code=404, detail="Person not found")

        # Get sources
        sources = (await session.execute(
            select(PersonSource).where(PersonSource.person_id == person_id)
            .order_by(PersonSource.relevance_score.desc())
        )).scalars().all()

        # Get jobs
        jobs = (await session.execute(
            select(DiscoveryJob).where(DiscoveryJob.person_id == person_id)
            .order_by(DiscoveryJob.created_at.desc())
        )).scalars().all()

        result = _person_to_dict(person)
        result["sources"] = [
            {
                "id": s.id,
                "source_type": s.source_type,
                "platform": s.platform,
                "url": s.url,
                "title": s.title,
                "raw_content": s.raw_content[:2000] if s.raw_content else None,
                "structured_data": json.loads(s.structured_data) if s.structured_data else None,
                "confidence": round(max(s.relevance_score or 0, s.source_reliability or 0), 2),
                "relevance_score": s.relevance_score,
                "source_reliability": s.source_reliability,
                "fetched_at": s.fetched_at.isoformat() if s.fetched_at else "",
            }
            for s in sources
        ]
        result["jobs"] = [
            {
                "id": j.id,
                "status": j.status,
                "total_cost": j.total_cost,
                "latency_ms": j.latency_ms,
                "sources_hit": j.sources_hit,
                "cache_hits": j.cache_hits,
                "created_at": j.created_at.isoformat() if j.created_at else "",
                "completed_at": j.completed_at.isoformat() if j.completed_at else None,
            }
            for j in jobs
        ]

        return result


@router.put("/persons/{person_id}")
async def update_person(person_id: str, update: PersonUpdate, _admin=Depends(require_admin)):
    factory = get_session_factory()
    async with factory() as session:
        person = (await session.execute(
            select(Person).where(Person.id == person_id)
        )).scalar_one_or_none()
        if not person:
            raise HTTPException(status_code=404, detail="Person not found")

        _JSON_FIELDS = {"key_facts", "education", "career_timeline", "expertise", "notable_work", "social_links",
                        "skills", "projects", "recommendations"}
        update_data = update.model_dump(exclude_none=True)
        for key, value in update_data.items():
            # JSON columns are stored as strings in the DB — serialize lists/dicts
            if key in _JSON_FIELDS and isinstance(value, (list, dict)):
                setattr(person, key, json.dumps(value))
            else:
                setattr(person, key, value)
        person.version += 1

        version = PersonVersion(
            person_id=person.id,
            version_number=person.version,
            profile_snapshot=json.dumps(_person_to_dict(person)),
            diff_from_previous=json.dumps(update_data),
            trigger="manual_edit",
        )
        session.add(version)
        await session.commit()

        result = _person_to_dict(person)
        sources = (await session.execute(
            select(PersonSource).where(PersonSource.person_id == person_id)
            .order_by(PersonSource.relevance_score.desc())
        )).scalars().all()
        jobs = (await session.execute(
            select(DiscoveryJob).where(DiscoveryJob.person_id == person_id)
            .order_by(DiscoveryJob.created_at.desc())
        )).scalars().all()
        result["sources"] = [
            {
                "id": s.id, "source_type": s.source_type, "platform": s.platform,
                "url": s.url, "title": s.title,
                "raw_content": s.raw_content[:2000] if s.raw_content else None,
                "structured_data": json.loads(s.structured_data) if s.structured_data else None,
                "confidence": round(max(s.relevance_score or 0, s.source_reliability or 0), 2),
                "relevance_score": s.relevance_score,
                "source_reliability": s.source_reliability,
                "fetched_at": s.fetched_at.isoformat() if s.fetched_at else "",
            }
            for s in sources
        ]
        result["jobs"] = [
            {
                "id": j.id, "status": j.status, "total_cost": j.total_cost,
                "latency_ms": j.latency_ms, "sources_hit": j.sources_hit,
                "cache_hits": j.cache_hits,
                "created_at": j.created_at.isoformat() if j.created_at else "",
                "completed_at": j.completed_at.isoformat() if j.completed_at else None,
            }
            for j in jobs
        ]
        return result


@router.delete("/persons/{person_id}")
async def delete_person(person_id: str, _admin=Depends(require_admin)):
    factory = get_session_factory()
    async with factory() as session:
        person = (await session.execute(
            select(Person).where(Person.id == person_id)
        )).scalar_one_or_none()
        if not person:
            raise HTTPException(status_code=404, detail="Person not found")

        # Delete sources, versions, jobs
        from sqlalchemy import delete as sql_delete
        await session.execute(sql_delete(PersonSource).where(PersonSource.person_id == person_id))
        await session.execute(sql_delete(PersonVersion).where(PersonVersion.person_id == person_id))
        await session.execute(sql_delete(DiscoveryJob).where(DiscoveryJob.person_id == person_id))
        await session.delete(person)
        await session.commit()

        return {"deleted": True}


@router.get("/persons/{person_id}/export")
async def export_person(
    person_id: str,
    format: str = Query("json", pattern="^(json|csv|pdf|pptx)$"),
    _auth=Depends(require_viewer),
):
    """Export person profile as JSON, CSV, or PDF."""
    _validate_uuid(person_id, "person_id")
    factory = get_session_factory()
    async with factory() as session:
        person = (await session.execute(
            select(Person).where(Person.id == person_id)
        )).scalar_one_or_none()
        if not person:
            raise HTTPException(status_code=404, detail="Person not found")

        sources = (await session.execute(
            select(PersonSource).where(PersonSource.person_id == person_id)
        )).scalars().all()

        profile = _person_to_dict(person)
        profile["sources"] = [
            {
                "source_type": s.source_type,
                "platform": s.platform,
                "url": s.url,
                "title": s.title,
                "confidence": round(max(s.relevance_score or 0, s.source_reliability or 0), 2),
            }
            for s in sources
        ]

    safe_name = profile.get("name", "export").replace(" ", "_")

    if format == "pdf":
        pdf_bytes = _generate_person_pdf(profile)
        return StreamingResponse(
            io.BytesIO(pdf_bytes),
            media_type="application/pdf",
            headers={"Content-Disposition": f'attachment; filename="{safe_name}_profile.pdf"'},
        )

    if format == "pptx":
        pptx_bytes = _generate_person_pptx(profile)
        return StreamingResponse(
            io.BytesIO(pptx_bytes),
            media_type="application/vnd.openxmlformats-officedocument.presentationml.presentation",
            headers={"Content-Disposition": f'attachment; filename="{safe_name}_profile.pptx"'},
        )

    if format == "csv":
        output = io.StringIO()
        writer = csv.writer(output)
        writer.writerow(["Field", "Value"])
        for key in ("name", "current_role", "company", "location", "bio", "confidence_score", "reputation_score"):
            writer.writerow([key, profile.get(key, "")])
        writer.writerow([])
        writer.writerow(["Source Platform", "URL", "Title", "Confidence"])
        for s in profile.get("sources", []):
            writer.writerow([s.get("platform", ""), s.get("url", ""), s.get("title", ""), s.get("confidence", "")])

        output.seek(0)
        return StreamingResponse(
            iter([output.getvalue()]),
            media_type="text/csv",
            headers={"Content-Disposition": f'attachment; filename="{safe_name}_profile.csv"'},
        )

    return JSONResponse(
        content=profile,
        headers={"Content-Disposition": f'attachment; filename="{safe_name}_profile.json"'},
    )


@router.post("/persons/{person_id}/re-search")
async def re_search_person(person_id: str, _admin=Depends(require_admin)):
    """Re-run discovery with the person's current data as input."""
    import asyncio

    factory = get_session_factory()
    async with factory() as session:
        person = (await session.execute(
            select(Person).where(Person.id == person_id)
        )).scalar_one_or_none()
        if not person:
            raise HTTPException(status_code=404, detail="Person not found")

        # Build input from existing person data
        social = person.get_json("social_links") or {}
        input_data = {
            "name": person.name,
            "company": person.company or "",
            "role": person.current_role or "",
            "location": person.location or "",
            "linkedin_url": social.get("linkedin", ""),
            "twitter_handle": social.get("twitter", ""),
            "github_username": social.get("github", ""),
            "instagram_handle": social.get("instagram", ""),
            "context": f"Re-search of existing profile. Previous bio: {(person.bio or '')[:200]}",
        }

        job_id = str(uuid.uuid4())
        job = DiscoveryJob(
            id=job_id,
            person_id=person.id,
            input_params=json.dumps(input_data),
            status="running",
        )
        session.add(job)
        await session.commit()

    asyncio.create_task(_run_discovery(job_id, input_data))

    return {"job_id": job_id, "status": "running", "message": "Re-search started."}


@router.post("/persons/{person_id}/refresh-image")
async def refresh_person_image(person_id: str, _admin=Depends(require_admin)):
    """
    Clear and re-resolve the profile image for a person without a full re-discovery.

    Useful when the stored image is incorrect (wrong person, landscape photo, etc.).
    The image resolver will run the full waterfall (LinkedIn, Wikipedia, Knowledge
    Graph, etc.) and store the best quality headshot it finds.
    """
    factory = get_session_factory()
    async with factory() as session:
        person = (await session.execute(
            select(Person).where(Person.id == person_id)
        )).scalar_one_or_none()
        if not person:
            raise HTTPException(status_code=404, detail="Person not found")

        name = person.name
        company = person.company

        # Get existing sources to give the resolver context (LinkedIn handle etc.)
        sources_rows = (await session.execute(
            select(PersonSource).where(PersonSource.person_id == person_id)
        )).scalars().all()
        search_results = []
        for s in sources_rows:
            entry: dict = {
                "url": s.url or "",
                "source_type": s.source_type or "",
            }
            if s.structured_data:
                try:
                    entry["structured"] = json.loads(s.structured_data)
                except Exception:
                    pass
            search_results.append(entry)

        # Clear old image so resolver doesn't hit the old (bad) cached value
        person.image_url = None
        await session.commit()

    # Run resolver outside the session
    from app.cache import set_cached_results  # noqa: PLC0415
    # Invalidate the old cache entry for this person
    old_cache_key = f"{name}|{company or ''}"
    for tool_name in ("image_resolver_v3", "image_resolver_v4"):
        try:
            await set_cached_results(old_cache_key, tool_name, [])
        except Exception:
            pass

    from app.tools.image_resolver import resolve_profile_image  # noqa: PLC0415
    new_image_url = await resolve_profile_image(name, company, search_results)

    if new_image_url:
        async with factory() as session:
            person = (await session.execute(
                select(Person).where(Person.id == person_id)
            )).scalar_one_or_none()
            if person:
                person.image_url = new_image_url
                await session.commit()
        return {
            "person_id": person_id,
            "name": name,
            "image_url": new_image_url,
            "message": "Image refreshed successfully.",
        }
    else:
        return {
            "person_id": person_id,
            "name": name,
            "image_url": None,
            "message": "No suitable image found. Profile left without image.",
        }


# --- Cost dashboard ---

@router.get("/admin/rate-limits")
async def get_rate_limits(_admin=Depends(require_admin)):
    """Get current per-source rate limit status."""
    return rate_limiter.get_status()


@router.get("/admin/costs")
async def get_cost_stats(_admin=Depends(require_admin)):
    factory = get_session_factory()
    async with factory() as session:
        # Total costs
        total = (await session.execute(
            select(func.sum(DiscoveryJob.total_cost))
        )).scalar() or 0.0

        # Total jobs
        job_count = (await session.execute(
            select(func.count()).select_from(DiscoveryJob)
        )).scalar() or 0

        # Avg cost
        avg_cost = (await session.execute(
            select(func.avg(DiscoveryJob.total_cost))
        )).scalar() or 0.0

        # Recent jobs with costs
        recent_jobs = (await session.execute(
            select(DiscoveryJob)
            .where(DiscoveryJob.status == "completed")
            .order_by(DiscoveryJob.created_at.desc())
            .limit(20)
        )).scalars().all()

        recent = []
        for j in recent_jobs:
            recent.append({
                "id": j.id,
                "total_cost": j.total_cost,
                "latency_ms": j.latency_ms,
                "sources_hit": j.sources_hit,
                "cache_hits": j.cache_hits,
                "created_at": j.created_at.isoformat() if j.created_at else "",
            })

        return {
            "total_spend": round(total, 4),
            "total_jobs": job_count,
            "average_cost": round(avg_cost, 4),
            "recent_jobs": recent,
        }


# --- Auth endpoint for admin login ---

class LoginRequest(BaseModel):
    email: str = Field("", max_length=255)
    password: str = Field("", max_length=128)


@router.post("/auth/login")
async def admin_login(request: LoginRequest):
    """Validate admin credentials and return a JWT token."""
    from app.auth import verify_admin, create_token_pair

    user = await verify_admin(request.email, request.password)
    if not user:
        raise HTTPException(status_code=401, detail="Invalid credentials")

    tokens = create_token_pair({"sub": user.email, "role": user.role})
    return {**tokens, "email": user.email, "role": user.role}


class RefreshRequest(BaseModel):
    refresh_token: str = Field(..., min_length=10)


class CreateUserRequest(BaseModel):
    email: str = Field(..., min_length=1, max_length=255)
    password: str = Field(..., min_length=8, max_length=128)
    role: str = Field("admin", pattern="^(admin|viewer|api_only)$")


@router.post("/auth/refresh")
async def refresh_token(request: RefreshRequest):
    from app.auth import verify_refresh_token, create_token_pair

    payload = verify_refresh_token(request.refresh_token)
    if not payload:
        raise HTTPException(status_code=401, detail="Invalid refresh token")

    tokens = create_token_pair({"sub": payload["sub"], "role": payload.get("role", "admin")})
    return tokens


# --- User management (admin only) ---

@router.post("/admin/users")
async def create_user(request: CreateUserRequest, _admin=Depends(require_admin)):
    """Create a new user. Admin only."""
    factory = get_session_factory()
    async with factory() as session:
        existing = (await session.execute(
            select(AdminUser).where(AdminUser.email == request.email)
        )).scalar_one_or_none()
        if existing:
            raise HTTPException(status_code=409, detail="User with this email already exists")

        user = AdminUser(
            email=request.email,
            password_hash=bcrypt.hash(request.password),
            role=request.role,
        )
        session.add(user)
        await session.commit()
        await session.refresh(user)
        return {
            "id": user.id,
            "email": user.email,
            "role": user.role,
            "created_at": user.created_at.isoformat() if user.created_at else "",
        }


@router.get("/admin/users")
async def list_users(_admin=Depends(require_admin)):
    """List all users. Admin only."""
    factory = get_session_factory()
    async with factory() as session:
        result = await session.execute(
            select(AdminUser).order_by(AdminUser.email)
        )
        users = result.scalars().all()
        return {
            "users": [
                {
                    "id": u.id,
                    "email": u.email,
                    "role": u.role,
                    "created_at": u.created_at.isoformat() if u.created_at else "",
                }
                for u in users
            ],
        }


@router.delete("/admin/users/{user_id}")
async def delete_user(user_id: str, _admin=Depends(require_admin)):
    """Delete a user. Admin only."""
    _validate_uuid(user_id, "user_id")
    factory = get_session_factory()
    async with factory() as session:
        user = (await session.execute(
            select(AdminUser).where(AdminUser.id == user_id)
        )).scalar_one_or_none()
        if not user:
            raise HTTPException(status_code=404, detail="User not found")
        await session.delete(user)
        await session.commit()
        return {"deleted": True}


# --- Utility ---

@router.get("/health")
async def health():
    checks = {"status": "healthy", "version": "2.0.0", "timestamp": time.time()}
    try:
        factory = get_session_factory()
        async with factory() as db:
            await db.execute(select(Person).limit(1))
        checks["database"] = "ok"
    except Exception:
        checks["database"] = "error"
        checks["status"] = "degraded"
    return checks


@router.post("/cache/cleanup")
async def cleanup_cache(_admin=Depends(require_admin)):
    count = await cleanup_expired_cache()
    return {"cleaned": count}


# --- Intelligence Endpoints ---

def _person_to_dict_for_intelligence(person: Person, sources: list[PersonSource]) -> dict:
    """Convert a Person model and its sources to a plain dict for intelligence analysis."""
    return {
        "name": person.name,
        "current_role": person.current_role,
        "company": person.company,
        "location": person.location,
        "bio": person.bio,
        "key_facts": person.get_json("key_facts") or [],
        "education": person.get_json("education") or [],
        "expertise": person.get_json("expertise") or [],
        "notable_work": person.get_json("notable_work") or [],
        "career_timeline": person.get_json("career_timeline") or [],
        "confidence_score": person.confidence_score,
        "reputation_score": person.reputation_score or 0,
        "social_links": person.get_json("social_links") or {},
        "sources": [
            {
                "title": s.title,
                "url": s.url,
                "platform": s.platform,
                "source_type": s.source_type,
                "raw_content": s.raw_content,
                "source_reliability": s.source_reliability,
                "relevance_score": s.relevance_score,
            }
            for s in sources
        ],
    }


@router.get("/persons/{person_id}/sentiment")
async def get_sentiment_analysis(person_id: str, _auth=Depends(require_viewer)):
    """Analyze public sentiment across all sources for a person."""
    factory = get_session_factory()
    async with factory() as session:
        person = (await session.execute(
            select(Person).where(Person.id == person_id)
        )).scalar_one_or_none()
        if not person:
            raise HTTPException(status_code=404, detail="Person not found")

        sources_result = await session.execute(
            select(PersonSource).where(PersonSource.person_id == person_id)
        )
        sources = sources_result.scalars().all()

        # Return cached sentiment if available and not requesting a refresh
        cached_sentiment = person.get_json("sentiment_data")
        if cached_sentiment and not isinstance(cached_sentiment, dict):
            cached_sentiment = None

        if cached_sentiment:
            return cached_sentiment

        profile = _person_to_dict_for_intelligence(person, sources)
        result = await analyze_sentiment(profile)

        # Persist sentiment result to DB for future fast retrieval
        if result and "error" not in result:
            try:
                person.set_json("sentiment_data", result)
                await session.commit()
            except Exception as e:
                logger.warning(f"Failed to cache sentiment for person {person_id}: {e}")

        return result


@router.get("/persons/{person_id}/influence")
async def get_influence_score(person_id: str, _auth=Depends(require_viewer)):
    """Calculate multi-dimensional influence score for a person."""
    factory = get_session_factory()
    async with factory() as session:
        person = (await session.execute(
            select(Person).where(Person.id == person_id)
        )).scalar_one_or_none()
        if not person:
            raise HTTPException(status_code=404, detail="Person not found")

        # Return cached influence score if already computed
        if person.influence_score is not None:
            return {
                "overall_influence_score": person.influence_score,
                "_cached": True,
                "_note": "Cached score. Re-fetch person to trigger recomputation.",
            }

        sources_result = await session.execute(
            select(PersonSource).where(PersonSource.person_id == person_id)
        )
        sources = sources_result.scalars().all()

        profile = _person_to_dict_for_intelligence(person, sources)
        result = await calculate_influence_score(profile)

        # Persist the overall influence score to DB
        if result and "error" not in result:
            raw_score = result.get("overall_influence_score")
            if raw_score is not None:
                try:
                    # Normalise to 0.0-1.0 range (API returns 0-100)
                    person.influence_score = round(float(raw_score) / 100.0, 4)
                    await session.commit()
                except Exception as e:
                    logger.warning(f"Failed to cache influence score for person {person_id}: {e}")

        return result


@router.post("/persons/relationships")
async def get_relationship_map(
    request: dict,
    _auth=Depends(require_viewer),
):
    """Map relationships between two persons."""
    person_a_id = request.get("person_a_id", "")
    person_b_id = request.get("person_b_id", "")
    if not person_a_id or not person_b_id:
        raise HTTPException(status_code=400, detail="person_a_id and person_b_id are required")

    factory = get_session_factory()
    async with factory() as session:
        person_a = (await session.execute(
            select(Person).where(Person.id == person_a_id)
        )).scalar_one_or_none()
        person_b = (await session.execute(
            select(Person).where(Person.id == person_b_id)
        )).scalar_one_or_none()

        if not person_a or not person_b:
            raise HTTPException(status_code=404, detail="One or both persons not found")

        sources_a = (await session.execute(
            select(PersonSource).where(PersonSource.person_id == person_a_id)
        )).scalars().all()
        sources_b = (await session.execute(
            select(PersonSource).where(PersonSource.person_id == person_b_id)
        )).scalars().all()

        profile_a = _person_to_dict_for_intelligence(person_a, sources_a)
        profile_b = _person_to_dict_for_intelligence(person_b, sources_b)
        result = await map_relationships(profile_a, profile_b)
        return result


@router.post("/persons/{person_id}/meeting-prep")
async def get_meeting_prep(
    person_id: str,
    request: dict | None = Body(None),
    _auth=Depends(require_viewer),
):
    """Generate AI-powered meeting preparation insights."""
    context = (request or {}).get("context", "")
    factory = get_session_factory()
    async with factory() as session:
        person = (await session.execute(
            select(Person).where(Person.id == person_id)
        )).scalar_one_or_none()
        if not person:
            raise HTTPException(status_code=404, detail="Person not found")

        sources_result = await session.execute(
            select(PersonSource).where(PersonSource.person_id == person_id)
        )
        sources = sources_result.scalars().all()

        profile = _person_to_dict_for_intelligence(person, sources)
        result = await generate_meeting_prep(profile, context)
        return result


@router.get("/persons/{person_id}/verify")
async def get_fact_verification(person_id: str, _auth=Depends(require_viewer)):
    """Cross-reference facts from multiple sources and flag inconsistencies."""
    factory = get_session_factory()
    async with factory() as session:
        person = (await session.execute(
            select(Person).where(Person.id == person_id)
        )).scalar_one_or_none()
        if not person:
            raise HTTPException(status_code=404, detail="Person not found")

        sources_result = await session.execute(
            select(PersonSource).where(PersonSource.person_id == person_id)
        )
        sources = sources_result.scalars().all()

        profile = _person_to_dict_for_intelligence(person, sources)
        result = await verify_facts(profile)
        return result


def _safe(text: str, max_chars: int = 0) -> str:
    """Encode text to Latin-1, replacing unsupported characters, for fpdf2 built-in fonts."""
    if not isinstance(text, str):
        text = str(text)
    # Replace common Unicode punctuation with ASCII equivalents
    replacements = {
        "\u2018": "'", "\u2019": "'", "\u201c": '"', "\u201d": '"',
        "\u2013": "-", "\u2014": "-", "\u2026": "...", "\u2022": "*",
        "\u00a0": " ", "\u00b7": "*",
    }
    for src, dst in replacements.items():
        text = text.replace(src, dst)
    # Drop anything still outside Latin-1
    text = text.encode("latin-1", errors="replace").decode("latin-1")
    if max_chars:
        text = text[:max_chars]
    return text


def _generate_person_pdf(profile: dict) -> bytes:
    """Generate a styled PDF report for a person profile."""
    from fpdf import FPDF

    pdf = FPDF()
    pdf.set_auto_page_break(auto=True, margin=20)
    pdf.add_page()

    # Title
    pdf.set_font("Helvetica", "B", 20)
    pdf.cell(0, 12, _safe(profile.get("name", "Unknown")), new_x="LMARGIN", new_y="NEXT")
    pdf.set_draw_color(41, 128, 185)
    pdf.line(10, pdf.get_y(), 200, pdf.get_y())
    pdf.ln(4)

    # Subtitle: role + company
    role_parts = [profile.get("current_role"), profile.get("company")]
    subtitle = " at ".join(p for p in role_parts if p)
    if subtitle:
        pdf.set_font("Helvetica", "I", 12)
        pdf.set_text_color(100, 100, 100)
        pdf.cell(0, 8, _safe(subtitle), new_x="LMARGIN", new_y="NEXT")
        pdf.set_text_color(0, 0, 0)

    # Metadata row
    meta_items = []
    if profile.get("location"):
        meta_items.append(f"Location: {_safe(profile['location'])}")
    meta_items.append(f"Confidence: {round((profile.get('confidence_score') or 0) * 100)}%")
    if profile.get("reputation_score") is not None:
        meta_items.append(f"Reputation: {round(profile['reputation_score'] * 100)}%")
    pdf.set_font("Helvetica", "", 9)
    pdf.set_text_color(120, 120, 120)
    pdf.cell(0, 6, "  |  ".join(meta_items), new_x="LMARGIN", new_y="NEXT")
    pdf.set_text_color(0, 0, 0)
    pdf.ln(4)

    def _section(title: str):
        pdf.set_font("Helvetica", "B", 13)
        pdf.set_text_color(41, 128, 185)
        pdf.cell(0, 10, _safe(title), new_x="LMARGIN", new_y="NEXT")
        pdf.set_text_color(0, 0, 0)

    def _body(text: str):
        pdf.set_font("Helvetica", "", 10)
        pdf.multi_cell(0, 5, _safe(text))
        pdf.ln(2)

    def _bullet_list(items: list):
        pdf.set_font("Helvetica", "", 10)
        for item in items:
            label = str(item) if isinstance(item, str) else json.dumps(item) if isinstance(item, dict) else str(item)
            # Prefix bullet inline — avoids cell+multi_cell x-position issues
            pdf.multi_cell(0, 5, "- " + _safe(label, max_chars=200))
            pdf.ln(1)

    if profile.get("bio"):
        _section("Bio")
        _body(profile["bio"])

    if profile.get("expertise"):
        _section("Expertise")
        _bullet_list(profile["expertise"])
        pdf.ln(2)

    if profile.get("key_facts"):
        _section("Key Facts")
        _bullet_list(profile["key_facts"])
        pdf.ln(2)

    if profile.get("education"):
        _section("Education")
        for edu in profile["education"]:
            if isinstance(edu, dict):
                parts = [edu.get("degree", ""), edu.get("institution", ""), edu.get("year", "")]
                _body(" - ".join(str(p) for p in parts if p))
            else:
                _body(str(edu))

    if profile.get("career_timeline"):
        _section("Career Timeline")
        for entry in profile["career_timeline"]:
            if isinstance(entry, dict):
                period = entry.get("period", entry.get("year", ""))
                role = entry.get("role", entry.get("title", ""))
                org = entry.get("company", entry.get("organization", ""))
                line = f"{period}: {role}"
                if org:
                    line += f" at {org}"
                _body(line)
            else:
                _body(str(entry))

    if profile.get("notable_work"):
        _section("Notable Work")
        _bullet_list(profile["notable_work"])
        pdf.ln(2)

    if profile.get("social_links"):
        _section("Social Links")
        pdf.set_font("Helvetica", "", 10)
        for platform, url in profile["social_links"].items():
            if url:
                pdf.cell(0, 5, _safe(f"{platform}: {url}"), new_x="LMARGIN", new_y="NEXT")
        pdf.ln(2)

    if profile.get("sources"):
        _section("Sources")
        pdf.set_font("Helvetica", "B", 9)
        pdf.cell(35, 6, "Platform")
        pdf.cell(100, 6, "Title")
        pdf.cell(25, 6, "Confidence", align="R")
        pdf.ln()
        pdf.set_font("Helvetica", "", 9)
        for s in profile["sources"][:30]:
            pdf.cell(35, 5, _safe(str(s.get("platform", "")), max_chars=20))
            pdf.cell(100, 5, _safe(str(s.get("title", "")), max_chars=60))
            pdf.cell(25, 5, f"{round((s.get('confidence', 0)) * 100)}%", align="R")
            pdf.ln()

    # Footer
    pdf.ln(6)
    pdf.set_font("Helvetica", "I", 8)
    pdf.set_text_color(150, 150, 150)
    pdf.cell(0, 5, _safe(f"Generated by People Discovery Agent | Version {profile.get('version', 1)}"), align="C")

    return pdf.output()


def _generate_person_pptx(profile: dict) -> bytes:
    """Generate a professional one-pager PPTX presentation for a person profile."""
    try:
        from pptx import Presentation
        from pptx.util import Inches, Pt
        from pptx.dml.color import RGBColor

        prs = Presentation()
        prs.slide_width = Inches(13.333)
        prs.slide_height = Inches(7.5)

        _BG = RGBColor(0x1A, 0x1A, 0x2E)
        _TITLE = RGBColor(255, 255, 255)
        _BODY = RGBColor(0xE0, 0xE0, 0xE0)
        _ACCENT = RGBColor(0x3B, 0x82, 0xF6)

        def _apply_dark_bg(slide):
            background = slide.background
            fill = background.fill
            fill.solid()
            fill.fore_color.rgb = _BG

        # Slide 1: Title
        slide_layout = prs.slide_layouts[6]
        s1 = prs.slides.add_slide(slide_layout)
        _apply_dark_bg(s1)
        title_box = s1.shapes.add_textbox(Inches(0.5), Inches(2.5), Inches(12.333), Inches(1.5))
        tf = title_box.text_frame
        p = tf.paragraphs[0]
        p.text = profile.get("name") or "Unknown"
        p.font.size = Pt(44)
        p.font.bold = True
        p.font.color.rgb = _TITLE
        p.alignment = 1

        role_parts = [profile.get("current_role"), profile.get("company")]
        subtitle = " at ".join(p for p in role_parts if p) or ""
        if subtitle:
            sub_box = s1.shapes.add_textbox(Inches(0.5), Inches(4.2), Inches(12.333), Inches(0.8))
            sub_tf = sub_box.text_frame
            sub_p = sub_tf.paragraphs[0]
            sub_p.text = subtitle
            sub_p.font.size = Pt(24)
            sub_p.font.color.rgb = _BODY
            sub_p.alignment = 1

        # Slide 2: Bio, Key Facts, Expertise
        s2 = prs.slides.add_slide(slide_layout)
        _apply_dark_bg(s2)
        left, top, width, height = Inches(0.5), Inches(0.5), Inches(12.333), Inches(6.5)
        content_box = s2.shapes.add_textbox(left, top, width, height)
        tf2 = content_box.text_frame
        tf2.word_wrap = True

        def _add_section(title: str, items: list | None = None, body: str | None = None):
            p = tf2.add_paragraph()
            p.text = title
            p.font.size = Pt(18)
            p.font.bold = True
            p.font.color.rgb = _ACCENT
            p.space_after = Pt(6)
            if body:
                bp = tf2.add_paragraph()
                bp.text = body[:1500] if body else ""
                bp.font.size = Pt(12)
                bp.font.color.rgb = _BODY
                bp.space_after = Pt(12)
            if items:
                for item in items[:20]:
                    label = str(item) if isinstance(item, str) else json.dumps(item)[:200] if isinstance(item, dict) else str(item)
                    ip = tf2.add_paragraph()
                    ip.text = f"• {label[:300]}"
                    ip.font.size = Pt(11)
                    ip.font.color.rgb = _BODY
                    ip.space_after = Pt(4)
                tf2.add_paragraph()

        if profile.get("bio"):
            _add_section("Bio", body=profile["bio"])
        if profile.get("key_facts"):
            _add_section("Key Facts", items=profile["key_facts"])
        if profile.get("expertise"):
            _add_section("Expertise", items=profile["expertise"])
        if not any([profile.get("bio"), profile.get("key_facts"), profile.get("expertise")]):
            _add_section("Overview", body="No additional details available.")

        # Slide 3: Career & Sources
        s3 = prs.slides.add_slide(slide_layout)
        _apply_dark_bg(s3)
        career_box = s3.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(12.333), Inches(6.5))
        tf3 = career_box.text_frame
        tf3.word_wrap = True

        p3 = tf3.paragraphs[0]
        p3.text = "Career & Sources"
        p3.font.size = Pt(18)
        p3.font.bold = True
        p3.font.color.rgb = _ACCENT
        p3.space_after = Pt(8)

        sources_count = len(profile.get("sources") or [])
        sp = tf3.add_paragraph()
        sp.text = f"Sources: {sources_count}"
        sp.font.size = Pt(12)
        sp.font.color.rgb = _BODY
        sp.space_after = Pt(12)

        for entry in (profile.get("career_timeline") or [])[:15]:
            if isinstance(entry, dict):
                period = entry.get("period", entry.get("year", ""))
                role = entry.get("role", entry.get("title", ""))
                org = entry.get("company", entry.get("organization", ""))
                line = f"{period}: {role}"
                if org:
                    line += f" at {org}"
            else:
                line = str(entry)
            ep = tf3.add_paragraph()
            ep.text = f"• {line[:200]}"
            ep.font.size = Pt(11)
            ep.font.color.rgb = _BODY
            ep.space_after = Pt(4)

        buf = io.BytesIO()
        prs.save(buf)
        buf.seek(0)
        return buf.getvalue()
    except Exception as e:
        logger.exception("PPTX generation failed: %s", e)
        raise HTTPException(status_code=500, detail=f"Failed to generate PPTX: {str(e)}")


def _person_to_dict(person: Person) -> dict:
    return {
        "id": person.id,
        "name": person.name,
        "current_role": person.current_role,
        "company": person.company,
        "location": person.location,
        "bio": person.bio,
        "image_url": person.image_url,
        "education": person.get_json("education"),
        "key_facts": person.get_json("key_facts"),
        "social_links": person.get_json("social_links"),
        "expertise": person.get_json("expertise"),
        "notable_work": person.get_json("notable_work"),
        "career_timeline": person.get_json("career_timeline"),
        "skills": person.get_json("skills"),
        "projects": person.get_json("projects"),
        "recommendations": person.get_json("recommendations"),
        "followers_count": person.followers_count,
        "blog_url": person.blog_url,
        "confidence_score": person.confidence_score,
        "reputation_score": person.reputation_score,
        "status": person.status,
        "version": person.version,
        "created_at": person.created_at.isoformat() if person.created_at else "",
        "updated_at": person.updated_at.isoformat() if person.updated_at else "",
    }


# ---------------------------------------------------------------------------
# Helpers for /fields and /summary
# ---------------------------------------------------------------------------

_SCALAR_FIELDS = {
    "name", "current_role", "company", "location", "bio", "image_url",
    "confidence_score", "reputation_score", "status", "version",
    "created_at", "updated_at",
}
_JSON_FIELDS = {
    "education", "key_facts", "social_links", "expertise",
    "notable_work", "career_timeline", "skills", "projects", "recommendations",
}
_SCALAR_ENRICHMENT_FIELDS = {"followers_count", "blog_url"}
_ALL_FIELDS = _SCALAR_FIELDS | _JSON_FIELDS | _SCALAR_ENRICHMENT_FIELDS

# Which platforms are most authoritative for each field.
# Sources matching these platforms are ranked higher in the per-field source list.
_FIELD_PLATFORM_AFFINITY: dict[str, list[str]] = {
    "name":            ["linkedin", "twitter", "wikipedia", "crunchbase"],
    "current_role":    ["linkedin", "crunchbase"],
    "company":         ["linkedin", "crunchbase"],
    "location":        ["linkedin", "twitter"],
    "bio":             ["linkedin", "wikipedia", "twitter"],
    "image_url":       ["linkedin", "twitter", "wikipedia"],
    "social_links":    ["linkedin", "twitter", "github"],
    "education":       ["linkedin", "wikipedia"],
    "career_timeline": ["linkedin", "crunchbase"],
    "notable_work":    ["wikipedia", "google_news", "web"],
    "expertise":       ["linkedin", "web", "google_news"],
    "key_facts":       ["linkedin", "wikipedia", "google_news", "web"],
}

# Source types that are always high-signal regardless of field
_HIGH_SIGNAL_SOURCE_TYPES = {"linkedin_profile", "wikipedia", "crunchbase"}


def _pick_fields(person: Person, fields: list[str]) -> dict:
    """Return only the requested fields from a Person ORM object."""
    out: dict = {}
    for f in fields:
        if f in _SCALAR_FIELDS:
            val = getattr(person, f, None)
            out[f] = val.isoformat() if hasattr(val, "isoformat") else val
        elif f in _JSON_FIELDS:
            out[f] = person.get_json(f)
    return out


def _source_to_dict(src: PersonSource) -> dict:
    """Serialise a PersonSource row to the public API shape."""
    rel = src.relevance_score or 0.0
    reliability = src.source_reliability or 0.0
    # Combined confidence = weighted average of relevance + reliability
    confidence = round((rel * 0.6 + reliability * 0.4), 3)
    return {
        "platform": src.platform,
        "source_type": src.source_type,
        "url": src.url,
        "title": src.title,
        "confidence_score": confidence,
        "relevance_score": round(rel, 3),
        "source_reliability": round(reliability, 3),
        "scorer_reason": src.scorer_reason or None,
        "fetched_at": src.fetched_at.isoformat() if src.fetched_at else None,
    }


def _sources_for_field(
    all_sources: list,
    field: str,
    limit: int = 5,
) -> list[dict]:
    """
    Return the top *limit* sources most relevant to *field*.

    Ranking logic (descending priority):
      1. Sources whose source_type is in _HIGH_SIGNAL_SOURCE_TYPES
      2. Sources whose platform matches _FIELD_PLATFORM_AFFINITY[field]
         (ranked by position in the affinity list — first = most relevant)
      3. Fallback: sorted by relevance_score desc
    Only sources with relevance_score ≥ 0.5 are included unless there
    aren't enough, in which case the threshold drops to 0.
    """
    affinity = _FIELD_PLATFORM_AFFINITY.get(field, [])

    def sort_key(s):
        high_signal = s.source_type in _HIGH_SIGNAL_SOURCE_TYPES
        affinity_rank = affinity.index(s.platform) if s.platform in affinity else len(affinity)
        rel = s.relevance_score or 0.0
        # Lower tuple → higher priority (we sort ascending then reverse)
        return (not high_signal, affinity_rank, -rel)

    candidates = [s for s in all_sources if (s.relevance_score or 0) >= 0.5]
    if len(candidates) < limit:
        candidates = list(all_sources)  # relax threshold

    ranked = sorted(candidates, key=sort_key)
    return [_source_to_dict(s) for s in ranked[:limit]]


# ---------------------------------------------------------------------------
# GET /persons/{person_id}/fields?fields=name,company,image_url,...
# ---------------------------------------------------------------------------
@router.get("/persons/{person_id}/fields")
async def get_person_fields(
    person_id: str,
    fields: str = "name,current_role,company,image_url",
    _auth=Depends(require_viewer),
):
    """
    Return **specific fields** for a person, each annotated with the top
    sources that support that field's value and a per-field confidence score.

    ### Query parameter
    - **fields** — comma-separated field names (default: `name,current_role,company,image_url`)

    ### Available fields
    `name`, `current_role`, `company`, `location`, `bio`, `image_url`,
    `education`, `key_facts`, `social_links`, `expertise`, `notable_work`,
    `career_timeline`, `confidence_score`, `reputation_score`,
    `status`, `version`, `created_at`, `updated_at`

    ### Response shape
    ```json
    {
      "id": "...",
      "person": { "name": "...", "company": "..." },
      "overall_confidence": 0.85,
      "fields": {
        "name": {
          "value": "Sam Altman",
          "confidence_score": 0.97,
          "sources": [
            {
              "platform": "linkedin",
              "url": "...",
              "title": "...",
              "confidence_score": 0.95,
              "relevance_score": 0.95,
              "source_reliability": 0.95,
              "scorer_reason": "LinkedIn profile for the exact target person",
              "fetched_at": "2026-03-13T..."
            }
          ]
        }
      }
    }
    ```

    ### Example
    ```
    GET /api/persons/{id}/fields?fields=name,company,current_role,location,bio,social_links
    ```
    """
    requested = [f.strip() for f in fields.split(",") if f.strip()]
    unknown = [f for f in requested if f not in _ALL_FIELDS]
    if unknown:
        raise HTTPException(
            status_code=422,
            detail=f"Unknown fields: {unknown}. Available: {sorted(_ALL_FIELDS)}",
        )

    factory = get_session_factory()
    async with factory() as session:
        person = (await session.execute(
            select(Person).where(Person.id == person_id)
        )).scalar_one_or_none()
        if not person:
            raise HTTPException(status_code=404, detail="Person not found")

        sources_result = await session.execute(
            select(PersonSource)
            .where(PersonSource.person_id == person_id)
            .order_by(PersonSource.relevance_score.desc())
        )
        all_sources = sources_result.scalars().all()

    field_values = _pick_fields(person, requested)

    # Build per-field blocks with field-appropriate source ranking
    fields_out: dict = {}
    for f in requested:
        # Skip meta/score fields from source annotation — they don't have sources
        if f in {"confidence_score", "reputation_score", "status", "version",
                  "created_at", "updated_at"}:
            fields_out[f] = {"value": field_values.get(f), "sources": []}
            continue

        field_sources = _sources_for_field(all_sources, f, limit=5)
        # Per-field confidence: average of top sources' confidence scores
        if field_sources:
            field_confidence = round(
                sum(s["confidence_score"] for s in field_sources) / len(field_sources), 3
            )
        else:
            field_confidence = person.confidence_score or 0.0

        fields_out[f] = {
            "value": field_values.get(f),
            "confidence_score": field_confidence,
            "sources": field_sources,
        }

    return {
        "id": person.id,
        "person": {
            "name": person.name,
            "company": person.company,
            "current_role": person.current_role,
        },
        "overall_confidence": person.confidence_score,
        "total_sources": len(all_sources),
        "fields": fields_out,
    }


# ---------------------------------------------------------------------------
# GET /persons/{person_id}/summary
# ---------------------------------------------------------------------------
@router.get("/persons/{person_id}/summary")
async def get_person_summary(
    person_id: str,
    _auth=Depends(require_viewer),
):
    """
    Return a **complete summary** of a person — all key profile fields with
    no raw source list. Ideal for cards, previews, and CRM sync.

    ### Response shape
    ```json
    {
      "id": "...",
      "name": "Sam Altman",
      "current_role": "CEO",
      "company": "OpenAI",
      "location": "San Francisco, CA",
      "image_url": "https://...",
      "bio": "Full biography text...",
      "expertise": ["AI", "Entrepreneurship"],
      "key_facts": ["Co-founded OpenAI in 2015", ...],
      "notable_work": [...],
      "education": [...],
      "career_timeline": [...],
      "social_links": { "linkedin": "...", "twitter": "..." },
      "confidence_score": 0.85,
      "reputation_score": null,
      "sources_count": 67,
      "last_updated": "2026-03-14T..."
    }
    ```
    """
    factory = get_session_factory()
    async with factory() as session:
        person = (await session.execute(
            select(Person).where(Person.id == person_id)
        )).scalar_one_or_none()
        if not person:
            raise HTTPException(status_code=404, detail="Person not found")

        from sqlalchemy import func as sa_func  # noqa: PLC0415
        sources_count_result = await session.execute(
            select(sa_func.count()).select_from(PersonSource)
            .where(PersonSource.person_id == person_id)
        )
        sources_count = sources_count_result.scalar() or 0

    return {
        "id": person.id,
        "name": person.name,
        "current_role": person.current_role,
        "company": person.company,
        "location": person.location,
        "image_url": person.image_url,
        "bio": person.bio,
        "expertise": person.get_json("expertise") or [],
        "key_facts": person.get_json("key_facts") or [],
        "notable_work": person.get_json("notable_work") or [],
        "education": person.get_json("education") or [],
        "career_timeline": person.get_json("career_timeline") or [],
        "social_links": person.get_json("social_links") or {},
        "skills": person.get_json("skills") or [],
        "projects": person.get_json("projects") or [],
        "recommendations": person.get_json("recommendations") or [],
        "followers_count": person.followers_count,
        "blog_url": person.blog_url,
        "confidence_score": person.confidence_score,
        "reputation_score": person.reputation_score,
        "sources_count": sources_count,
        "last_updated": person.updated_at.isoformat() if person.updated_at else None,
    }


# ── Staleness / Auto-Refresh ─────────────────────────────────────────────────

@router.get("/admin/staleness")
async def get_staleness_status(_auth=Depends(require_admin)):
    """Return how many persons are stale and when they were last updated."""
    from app.staleness import STALE_AFTER_DAYS, REFRESH_COOLDOWN_DAYS, BATCH_SIZE, CRON_INTERVAL_SECS
    from datetime import timedelta

    factory = get_session_factory()
    async with factory() as session:
        stale_cutoff = datetime.now(timezone.utc) - timedelta(days=STALE_AFTER_DAYS)
        total_q = select(func.count()).select_from(Person)
        stale_q = select(func.count()).select_from(Person).where(Person.updated_at < stale_cutoff)
        total = (await session.execute(total_q)).scalar() or 0
        stale = (await session.execute(stale_q)).scalar() or 0

        oldest = (await session.execute(
            select(Person.name, Person.updated_at).order_by(Person.updated_at.asc()).limit(5)
        )).all()

    return {
        "total_persons": total,
        "stale_count": stale,
        "stale_after_days": STALE_AFTER_DAYS,
        "refresh_cooldown_days": REFRESH_COOLDOWN_DAYS,
        "batch_size": BATCH_SIZE,
        "cron_interval_seconds": CRON_INTERVAL_SECS,
        "oldest_profiles": [
            {"name": r.name, "updated_at": r.updated_at.isoformat() if r.updated_at else None}
            for r in oldest
        ],
    }


@router.post("/admin/staleness/trigger")
async def trigger_staleness_refresh(_auth=Depends(require_admin)):
    """Manually trigger one staleness cron tick (for testing / admin use)."""
    from app.staleness import staleness_cron_tick
    asyncio.create_task(staleness_cron_tick(), name="manual-stale-tick")
    return {"message": "Staleness refresh tick queued"}


